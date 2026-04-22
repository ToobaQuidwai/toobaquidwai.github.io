#!/usr/bin/env python3
from __future__ import annotations

import math
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

import imageio.v2 as imageio
import numpy as np
from PIL import Image, ImageChops, ImageDraw, ImageEnhance, ImageFilter, ImageFont, ImageOps
from pptx import Presentation


EMU_PER_INCH = 914400
FRAME_SIZE = (1920, 1088)
FPS = 15


@dataclass
class SlideAssets:
    index: int
    main_image: Image.Image
    inset_image: Image.Image | None
    roi: tuple[float, float, float, float] | None


def ease_in_out(t: float) -> float:
    return 0.5 - 0.5 * math.cos(math.pi * max(0.0, min(1.0, t)))


def lerp(a: float, b: float, t: float) -> float:
    return a + (b - a) * t


def mix_box(
    start: tuple[float, float, float, float],
    end: tuple[float, float, float, float],
    t: float,
) -> tuple[float, float, float, float]:
    return tuple(lerp(a, b, t) for a, b in zip(start, end))


def center_crop_box(
    img_size: tuple[int, int],
    center: tuple[float, float],
    zoom: float = 1.0,
    target_ratio: float = FRAME_SIZE[0] / FRAME_SIZE[1],
) -> tuple[float, float, float, float]:
    img_w, img_h = img_size
    base_w = img_w / zoom
    base_h = base_w / target_ratio
    if base_h > img_h:
        base_h = img_h / zoom
        base_w = base_h * target_ratio
    cx, cy = center
    left = max(0.0, min(img_w - base_w, cx - base_w / 2))
    top = max(0.0, min(img_h - base_h, cy - base_h / 2))
    return (left, top, left + base_w, top + base_h)


def clean_image(img: Image.Image) -> Image.Image:
    img = ImageOps.exif_transpose(img).convert("RGB")
    w, h = img.size

    # Confocal screen captures in this deck include software chrome.
    if w / h > 2.2 and h <= 1800:
        left = int(w * 0.06)
        top = int(h * 0.06)
        right = int(w * 0.93)
        bottom = int(h * 0.82)
        cropped = img.crop((left, top, right, bottom))
    else:
        cropped = img

    cropped = ImageOps.autocontrast(cropped, cutoff=0.5)
    cropped = ImageEnhance.Contrast(cropped).enhance(1.18)
    cropped = ImageEnhance.Color(cropped).enhance(1.08)
    cropped = ImageEnhance.Sharpness(cropped).enhance(1.15)
    return cropped


def to_pil(blob: bytes) -> Image.Image:
    return Image.open(BytesIO(blob)).convert("RGB")


def extract_slide_assets(ppt_path: Path) -> list[SlideAssets]:
    prs = Presentation(str(ppt_path))
    slides: list[SlideAssets] = []

    for idx, slide in enumerate(prs.slides, start=1):
        pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
        pictures.sort(key=lambda s: s.width * s.height, reverse=True)
        rects = [shape for shape in slide.shapes if shape.shape_type == 1]

        if not pictures:
            continue

        main_shape = pictures[0]
        main_img = clean_image(to_pil(main_shape.image.blob))

        inset_img = None
        if len(pictures) > 1:
            inset_img = clean_image(to_pil(pictures[1].image.blob))

        roi = None
        if rects:
            rect = rects[0]
            rel_left = (rect.left - main_shape.left) / main_shape.width
            rel_top = (rect.top - main_shape.top) / main_shape.height
            rel_w = rect.width / main_shape.width
            rel_h = rect.height / main_shape.height

            img_w, img_h = main_img.size
            roi = (
                max(0.0, rel_left * img_w),
                max(0.0, rel_top * img_h),
                min(img_w, (rel_left + rel_w) * img_w),
                min(img_h, (rel_top + rel_h) * img_h),
            )

        slides.append(
            SlideAssets(
                index=idx,
                main_image=main_img,
                inset_image=inset_img,
                roi=roi,
            )
        )

    return slides


def fit_cover(img: Image.Image, size: tuple[int, int]) -> Image.Image:
    src_ratio = img.width / img.height
    dst_ratio = size[0] / size[1]
    if src_ratio > dst_ratio:
        crop_h = img.height
        crop_w = int(crop_h * dst_ratio)
    else:
        crop_w = img.width
        crop_h = int(crop_w / dst_ratio)
    left = (img.width - crop_w) // 2
    top = (img.height - crop_h) // 2
    return img.crop((left, top, left + crop_w, top + crop_h)).resize(size, Image.LANCZOS)


def render_cropped(img: Image.Image, box: tuple[float, float, float, float]) -> Image.Image:
    cropped = img.crop(tuple(int(v) for v in box))
    return fit_cover(cropped, FRAME_SIZE)


def add_vignette(img: Image.Image) -> Image.Image:
    w, h = img.size
    mask = Image.new("L", img.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse(
        [int(w * 0.04), int(h * 0.02), int(w * 0.96), int(h * 0.98)],
        fill=180,
    )
    mask = mask.filter(ImageFilter.GaussianBlur(90))
    shade = Image.new("RGB", img.size, (0, 0, 0))
    return Image.composite(img, Image.blend(img, shade, 0.28), mask)


def overlay_roi(frame: Image.Image, roi_box: tuple[float, float, float, float], full_box: tuple[float, float, float, float], alpha: float) -> Image.Image:
    if alpha <= 0:
        return frame

    overlay = frame.convert("RGBA")
    draw = ImageDraw.Draw(overlay)
    sx = FRAME_SIZE[0] / (full_box[2] - full_box[0])
    sy = FRAME_SIZE[1] / (full_box[3] - full_box[1])
    x1 = int((roi_box[0] - full_box[0]) * sx)
    y1 = int((roi_box[1] - full_box[1]) * sy)
    x2 = int((roi_box[2] - full_box[0]) * sx)
    y2 = int((roi_box[3] - full_box[1]) * sy)
    x1, x2 = sorted((max(0, min(FRAME_SIZE[0], x1)), max(0, min(FRAME_SIZE[0], x2))))
    y1, y2 = sorted((max(0, min(FRAME_SIZE[1], y1)), max(0, min(FRAME_SIZE[1], y2))))
    if x2 - x1 < 8 or y2 - y1 < 8:
        return frame

    stroke = (245, 250, 255, int(230 * alpha))
    for width in (10, 6, 3):
        draw.rounded_rectangle([x1, y1, x2, y2], radius=10, outline=stroke, width=width)
    return overlay.convert("RGB")


def build_frames(slides: list[SlideAssets]) -> list[Image.Image]:
    frames: list[Image.Image] = []

    for slide in slides:
        main = add_vignette(slide.main_image)
        full_box = center_crop_box(main.size, (main.width / 2, main.height / 2), zoom=1.0)

        if slide.roi:
            roi_w = max(120.0, slide.roi[2] - slide.roi[0])
            roi_h = max(120.0, slide.roi[3] - slide.roi[1])
            roi_center = ((slide.roi[0] + slide.roi[2]) / 2, (slide.roi[1] + slide.roi[3]) / 2)
            focus_zoom = min(
                7.0,
                max(
                    2.6,
                    min(main.width / (roi_w * 2.6), main.height / (roi_h * 2.6)),
                ),
            )
            focus_box = center_crop_box(main.size, roi_center, zoom=focus_zoom)
        else:
            focus_box = center_crop_box(main.size, (main.width / 2, main.height / 2), zoom=2.0)

        # Overview hold.
        for frame_idx in range(int(FPS * 0.8)):
            t = frame_idx / max(1, int(FPS * 0.8) - 1)
            drift_box = mix_box(full_box, center_crop_box(main.size, (main.width * 0.52, main.height * 0.48), zoom=1.06), ease_in_out(t))
            frame = render_cropped(main, drift_box)
            if slide.roi:
                frame = overlay_roi(frame, slide.roi, drift_box, alpha=min(1.0, t * 1.6))
            frames.append(frame)

        # Travel into the selected region.
        for frame_idx in range(int(FPS * 1.1)):
            t = ease_in_out(frame_idx / max(1, int(FPS * 1.1) - 1))
            box = mix_box(full_box, focus_box, t)
            frame = render_cropped(main, box)
            if slide.roi:
                frame = overlay_roi(frame, slide.roi, box, alpha=max(0.2, 1.0 - t * 0.4))
            frames.append(frame)

        if slide.inset_image:
            inset = add_vignette(slide.inset_image)
            inset_full = fit_cover(inset, FRAME_SIZE)
            main_focus = render_cropped(main, focus_box)

            # Crossfade from source region to the true zoom view.
            for frame_idx in range(int(FPS * 0.9)):
                t = ease_in_out(frame_idx / max(1, int(FPS * 0.9) - 1))
                frame = Image.blend(main_focus, inset_full, t)
                frames.append(frame)

            # Hold the zoomed view.
            for _ in range(int(FPS * 0.7)):
                frames.append(inset_full)

    # Finish on the final high-magnification view if available.
    if slides:
        last = slides[-1]
        end_img = fit_cover(add_vignette(last.inset_image or last.main_image), FRAME_SIZE)
        for _ in range(int(FPS * 1.2)):
            frames.append(end_img)

    return frames


def save_video(frames: list[Image.Image], output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    writer = imageio.get_writer(str(output_path), fps=FPS, codec="libx264", quality=8, pixelformat="yuv420p")
    try:
        for frame in frames:
            writer.append_data(np.asarray(frame))
    finally:
        writer.close()


def main() -> None:
    ppt_path = Path("/Users/toobaquidwai/Desktop/1.pptx")
    output_path = Path("/Users/toobaquidwai/Desktop/kidney_confocal_timelapse.mp4")
    slides = extract_slide_assets(ppt_path)
    frames = build_frames(slides)
    save_video(frames, output_path)
    print(output_path)


if __name__ == "__main__":
    main()
